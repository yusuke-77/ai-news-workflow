"""
Step 3: Excel のデータを読み込み PowerPoint を生成する
毎回同じデザインテンプレートで上書きするため、
「フォーマットを維持したまま内容を反映」できる。
"""
import os
from datetime import datetime

import openpyxl
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

BASE_DIR   = os.path.dirname(__file__)
EXCEL_PATH = os.path.join(BASE_DIR, "output", "ai_news.xlsx")
PPTX_PATH  = os.path.join(BASE_DIR, "output", "ai_news_presentation.pptx")

# ── カラーパレット ──────────────────────────────
C_DARK_BLUE   = RGBColor(0x1F, 0x38, 0x64)
C_MID_BLUE    = RGBColor(0x2E, 0x75, 0xB6)
C_LIGHT_BLUE  = RGBColor(0xBD, 0xD7, 0xEE)
C_PALE        = RGBColor(0xEE, 0xF3, 0xFB)
C_WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
C_ORANGE      = RGBColor(0xFF, 0x6B, 0x35)
C_TEXT        = RGBColor(0x1A, 0x1A, 0x2E)
C_MUTED       = RGBColor(0x88, 0x88, 0x99)
C_LINK        = RGBColor(0x44, 0x72, 0xC4)


# ── ユーティリティ ──────────────────────────────
def _solid_shape(slide, left, top, width, height, color: RGBColor, line=False):
    shp = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if line:
        shp.line.color.rgb = RGBColor(0xCC, 0xD9, 0xEA)
    else:
        shp.line.fill.background()
    return shp


def _textbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))


def _para(tf, text, size, bold=False, color=C_TEXT, align=PP_ALIGN.LEFT):
    p = tf.add_paragraph()
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return p


def _blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _bg(slide, color: RGBColor):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


# ── スライド: タイトル ──────────────────────────
def _title_slide(prs, date_str: str):
    slide = _blank_slide(prs)
    _bg(slide, C_DARK_BLUE)

    # アクセントライン
    _solid_shape(slide, 0, 3.55, 10, 0.06, C_ORANGE)

    # メインタイトル
    tb = _textbox(slide, 1.0, 1.4, 8.0, 1.3)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "最新 AI ニュース"
    r.font.size = Pt(46)
    r.font.bold = True
    r.font.color.rgb = C_WHITE

    # サブタイトル
    tb2 = _textbox(slide, 1.0, 3.0, 8.0, 0.65)
    tf2 = tb2.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = f"Top 10 Stories  |  {date_str}"
    r2.font.size = Pt(19)
    r2.font.color.rgb = C_LIGHT_BLUE

    # フッター
    tb3 = _textbox(slide, 1.0, 6.6, 8.0, 0.4)
    tf3 = tb3.text_frame
    p3 = tf3.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    r3 = p3.add_run()
    r3.text = "自動収集 by Claude Code"
    r3.font.size = Pt(11)
    r3.font.color.rgb = RGBColor(0x8A, 0xB4, 0xD4)


# ── スライド: 目次 ──────────────────────────────
def _index_slide(prs, news_items: list[dict]):
    slide = _blank_slide(prs)
    _bg(slide, C_PALE)
    _solid_shape(slide, 0, 0, 10, 1.1, C_DARK_BLUE)

    tb = _textbox(slide, 0.35, 0.2, 9.0, 0.7)
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "目次 ― Top 10 AI ニュース"
    r.font.size = Pt(22)
    r.font.bold = True
    r.font.color.rgb = C_WHITE

    for i, item in enumerate(news_items):
        y = 1.28 + i * 0.52
        # 番号バッジ
        badge = _solid_shape(slide, 0.28, y, 0.40, 0.38, C_MID_BLUE)
        btf = badge.text_frame
        btf.paragraphs[0].alignment = PP_ALIGN.CENTER
        br = btf.paragraphs[0].add_run()
        br.text = str(i + 1)
        br.font.size = Pt(12)
        br.font.bold = True
        br.font.color.rgb = C_WHITE

        # タイトル
        title = item["title"]
        if len(title) > 72:
            title = title[:69] + "…"
        tb2 = _textbox(slide, 0.80, y + 0.02, 8.4, 0.38)
        p2 = tb2.text_frame.paragraphs[0]
        r2 = p2.add_run()
        r2.text = title
        r2.font.size = Pt(12)
        r2.font.color.rgb = C_TEXT

        # ソースラベル（右端）
        tb3 = _textbox(slide, 8.7, y, 1.0, 0.36)
        p3 = tb3.text_frame.paragraphs[0]
        p3.alignment = PP_ALIGN.RIGHT
        r3 = p3.add_run()
        r3.text = item["source"][:16]
        r3.font.size = Pt(9)
        r3.font.color.rgb = C_MID_BLUE


# ── スライド: ニュース個別 ──────────────────────
def _news_slide(prs, item: dict, slide_num: int):
    slide = _blank_slide(prs)
    _bg(slide, C_WHITE)
    _solid_shape(slide, 0, 0, 10, 1.35, C_DARK_BLUE)

    # 番号バッジ
    badge = _solid_shape(slide, 0.25, 0.2, 0.58, 0.58, C_ORANGE)
    btf = badge.text_frame
    btf.paragraphs[0].alignment = PP_ALIGN.CENTER
    br = btf.paragraphs[0].add_run()
    br.text = str(slide_num)
    br.font.size = Pt(15)
    br.font.bold = True
    br.font.color.rgb = C_WHITE

    # タイトル
    title = item["title"] if len(item["title"]) <= 85 else item["title"][:82] + "…"
    tb = _textbox(slide, 0.98, 0.13, 8.75, 1.1)
    tb.text_frame.word_wrap = True
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(18)
    r.font.bold = True
    r.font.color.rgb = C_WHITE

    # 概要ボックス
    _solid_shape(slide, 0.35, 1.5, 9.3, 3.85, C_PALE, line=True)
    tb2 = _textbox(slide, 0.55, 1.6, 9.0, 3.65)
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    _para(tf2, "概要", 13, bold=True, color=C_MID_BLUE)

    p_body = tf2.add_paragraph()
    p_body.space_before = Pt(4)
    r_body = p_body.add_run()
    r_body.text = item["summary"] if item["summary"] else "（概要なし）"
    r_body.font.size = Pt(13)
    r_body.font.color.rgb = C_TEXT

    # メタ情報
    tb3 = _textbox(slide, 0.35, 5.5, 4.5, 0.42)
    p3 = tb3.text_frame.paragraphs[0]
    r3 = p3.add_run()
    r3.text = f"ソース: {item['source']}"
    r3.font.size = Pt(12)
    r3.font.color.rgb = C_MID_BLUE

    tb4 = _textbox(slide, 4.9, 5.5, 4.8, 0.42)
    p4 = tb4.text_frame.paragraphs[0]
    r4 = p4.add_run()
    r4.text = f"公開: {item['published']}"
    r4.font.size = Pt(12)
    r4.font.color.rgb = C_MUTED

    url_text = item["url"] if len(item["url"]) <= 95 else item["url"][:92] + "…"
    tb5 = _textbox(slide, 0.35, 6.05, 9.3, 0.42)
    p5 = tb5.text_frame.paragraphs[0]
    r5 = p5.add_run()
    r5.text = f"URL: {url_text}"
    r5.font.size = Pt(10)
    r5.font.color.rgb = C_LINK


# ── メイン ─────────────────────────────────────
def read_excel_data() -> list[dict]:
    wb = openpyxl.load_workbook(EXCEL_PATH)
    ws = wb["AIニュース"]
    rows = []
    for row in ws.iter_rows(min_row=2, values_only=True):
        if row[1]:
            rows.append(
                {
                    "title":     row[1],
                    "summary":   row[2] or "",
                    "url":       row[3] or "",
                    "published": row[4] or "",
                    "source":    row[5] or "",
                }
            )
    return rows


def create_pptx(news_items: list[dict] | None = None) -> str:
    os.makedirs(os.path.dirname(PPTX_PATH), exist_ok=True)

    if news_items is None:
        news_items = read_excel_data()

    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(7.5)

    date_str = datetime.now().strftime("%Y年%m月%d日")
    _title_slide(prs, date_str)
    _index_slide(prs, news_items)
    for i, item in enumerate(news_items, start=1):
        _news_slide(prs, item, i)

    prs.save(PPTX_PATH)
    print(f"  PowerPoint 生成完了: {PPTX_PATH}")
    return PPTX_PATH


if __name__ == "__main__":
    create_pptx()
